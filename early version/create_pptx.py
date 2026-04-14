from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
TEAL_DARK = RgbColor(0x1B, 0x3F, 0x1C)
TEAL = RgbColor(0x2E, 0x7D, 0x32)
TEAL_LIGHT = RgbColor(0xE8, 0xF5, 0xE9)
AMBER = RgbColor(0xB8, 0x86, 0x0B)
CORAL = RgbColor(0xA4, 0x58, 0x2A)
GRAY = RgbColor(0x8A, 0x7D, 0x6B)
HEADING = RgbColor(0x1B, 0x2A, 0x1C)
BODY = RgbColor(0x4A, 0x45, 0x39)
BG = RgbColor(0xFA, 0xFA, 0xF5)

def add_title_slide(prs, title, subtitle, meta=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark green background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = TEAL_DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(11.8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)
    p.alignment = PP_ALIGN.CENTER
    
    # Meta
    if meta:
        meta_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.3), Inches(1))
        tf = meta_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = meta
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, tag, title, desc=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Light green background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = TEAL_LIGHT
    background.line.fill.background()
    
    # Tag
    tag_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.8), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.8), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = HEADING
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    if desc:
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = BODY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, tag, title, content, tag_color=TEAL):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG
    background.line.fill.background()
    
    # Tag
    tag_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(3), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = tag_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(11.8), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = HEADING
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(11.8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = BODY
        p.space_after = Pt(12)
    
    return slide

def add_stat_slide(prs, tag, title, intro, stats, insight="", tag_color=TEAL):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG
    background.line.fill.background()
    
    # Tag
    tag_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(3), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = tag_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(11.8), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = HEADING
    
    # Intro
    intro_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(0.8))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = intro
    p.font.size = Pt(18)
    p.font.color.rgb = BODY
    
    # Stats
    stat_width = 3.5
    start_x = (13.333 - (stat_width * len(stats) + 0.3 * (len(stats) - 1))) / 2
    
    for i, (num, label, color) in enumerate(stats):
        x = start_x + i * (stat_width + 0.3)
        
        # Stat box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.8), Inches(stat_width), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(255, 255, 255)
        box.line.color.rgb = RgbColor(0xE0, 0xDD, 0xD4)
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(3), Inches(stat_width), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(stat_width - 0.4), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Insight
    if insight:
        insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.2), Inches(11.8), Inches(1.2))
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = TEAL_LIGHT
        insight_box.line.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)
        
        text_box = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(0.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insight
        p.font.size = Pt(16)
        p.font.color.rgb = TEAL_DARK
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, 
    "Designing Innovation:\nIncentives for Methane Reduction",
    "An Economic Design Analysis of an Advance Market Commitment for Methane Reduction",
    "Delany Broome | UC Berkeley, Master of Climate Solutions 2026\nIn collaboration with The Nature Conservancy & Spark Climate Solutions"
)

# SLIDE 2: The Problem
add_stat_slide(prs, "Situation", 
    "Methane is fast-acting, high-impact, and addressable",
    "Methane accounts for roughly 30% of global warming since pre-industrial times. But unlike CO2, it only persists in the atmosphere for about 12 years.",
    [("30%", "of global warming from methane", TEAL),
     ("60%", "of methane from human activities", AMBER),
     ("12 yrs", "atmospheric lifespan (vs 600+ for CO2)", CORAL)],
    "Key insight: Fast livestock methane cuts can deliver climate gains while longer-term decarbonization scales.",
    AMBER
)

# SLIDE 3: The Funnel
add_content_slide(prs, "Situation",
    "From global emissions to a single intervention",
    ["The emissions funnel narrows from all greenhouse gases down to a specific, addressable opportunity:",
     "",
     "Global greenhouse gas emissions",
     "    ↓",
     "Methane (CH4) emissions",
     "    ↓", 
     "Agriculture methane",
     "    ↓",
     "Enteric fermentation",
     "    ↓",
     "Cattle (75% of enteric CH4) ← AMC Target",
     "",
     "Sources: IPCC AR6, FAO GLEAM, EPA GHG Inventory 2024"],
    AMBER
)

# SLIDE 4: Economic Case
add_stat_slide(prs, "Situation",
    "The economic case for action is clear",
    "A single cow emits 154-264 lbs of methane per year, generating over $100 in annual climate damages.",
    [("$10B+", "Annual social cost of U.S. cattle methane (2022)", TEAL),
     ("80x", "Methane warming power vs CO2 over 20 years", AMBER)],
    "With U.S. cattle producing 6.64 million tons of methane in 2022, the opportunity for intervention is substantial.",
    AMBER
)

# SLIDE 5: The Deadlock
add_content_slide(prs, "Complication",
    "The innovation deadlock",
    ["Despite the clear economic case, demand uncertainty stalls private investment.",
     "",
     "🏢 Corporate Buyers ask:",
     "    \"How can I buy something that doesn't exist?\"",
     "",
     "🔬 Innovators ask:",
     "    \"Why should I build something no one is buying?\"",
     "",
     "Without reliable offtake, innovation stalls and capital arrives too late."],
    AMBER
)

# SLIDE 6: AMC Solution
add_content_slide(prs, "Resolution",
    "An AMC breaks the cycle",
    ["An Advance Market Commitment disrupts the deadlock by inserting a credible demand signal.",
     "",
     "Sponsors → AMC (Demand signal) → Innovators → Producers → Methane Reduced",
     "",
     "• Zero downside for sponsors: Pay nothing if no firm succeeds",
     "• Credible signal: Legally binding commitments create bankable demand",
     "• Capital unlocked: Investors can underwrite commercialization with clearer protection"],
    TEAL
)

# SLIDE 7: Mechanism Design
add_content_slide(prs, "Mechanism Design",
    "How the AMC mechanism works",
    ["Three payment pathways for sponsors:",
     "",
     "1. Offtake Premium",
     "   Pay per pound of low-methane beef, passed through supply chain",
     "",
     "2. Intervention Subsidy", 
     "   Pay per dose administered, covering vaccination costs",
     "",
     "3. Carbon Payment",
     "   Pay per ton of verified CO2e reduced",
     "",
     "The AMC is technology-agnostic: vaccines, genetics, feed additives, and slow-release boluses all qualify."],
    TEAL
)

# SLIDE 8: Technology Pathways
add_content_slide(prs, "Technology",
    "Technology pathways and readiness",
    ["🧬 Methane Vaccine (Early stage)",
     "   Immunological approach targeting rumen methanogens. Works across all systems.",
     "",
     "🧪 Feed Additives / 3-NOP (Commercializing)",
     "   Chemical inhibitor, proven ~30% reduction. Requires controlled feeding.",
     "",
     "💊 Slow-Release Bolus (In development)",
     "   Oral device releasing inhibitor over months.",
     "",
     "🐄 Genetic Selection (Long-term)",
     "   Breeding for naturally lower-emitting cattle."],
    TEAL
)

# SLIDE 9: Frameworks Section
add_section_slide(prs, "Analytical Frameworks",
    "MCS Program Frameworks",
    "Rigorous analytical methods applied to evaluate the methane vaccine advance market commitment"
)

# SLIDE 10: TEA Framework
add_stat_slide(prs, "Framework 1",
    "Techno-Economic Analysis: The MCS Playbook",
    "The TEA framework structures climate solution evaluation into three essential steps.",
    [("14.4x", "Benefit-Cost Ratio (baseline)", TEAL),
     ("$10.7B", "Net Present Value", TEAL)],
    "Step 1: Technical potential ($11.48B social benefits) | Step 2: Cost per ton CO2e | Step 3: Sensitivity analysis",
    RgbColor(0x7F, 0x77, 0xDD)
)

# SLIDE 11: LCA Framework
add_content_slide(prs, "Framework 2",
    "Lifecycle Assessment",
    ["Evaluating whether lifecycle emissions from vaccine production are justified by abatement.",
     "",
     "System Boundary:",
     "R&D inputs → Manufacturing → Distribution → Administration → Use phase (Methane reduction)",
     "",
     "Key Finding:",
     "The ratio of avoided emissions to embodied emissions is 100:1 or greater.",
     "",
     "Lifecycle emissions of vaccine production are negligible relative to abatement achieved.",
     "",
     "Methodology: ISO 14040/14044 framework. Functional unit: one vaccination course per animal."],
    RgbColor(0x7F, 0x77, 0xDD)
)

# SLIDE 12: Pathways & Barriers 1
add_content_slide(prs, "Framework 3",
    "Pathways and Barriers",
    ["The AMC functions as a niche-stimulation mechanism in transition theory terms.",
     "",
     "Barriers addressed by AMC:",
     "",
     "✓ Demand Uncertainty — Addressed",
     "   AMC inserts legally binding demand before commercialization",
     "",
     "◐ Grazing System Exclusion — Partially addressed",
     "   Vaccine works across all production systems",
     "",
     "◐ Adoption Friction — Partially addressed",
     "   AMC subsidizes vaccination, bundles with existing vaccines",
     "",
     "✓ Verification Complexity — Addressed",
     "   Uses vaccination verification rather than direct measurement"],
    RgbColor(0x7F, 0x77, 0xDD)
)

# SLIDE 13: Pathways & Barriers 2
add_content_slide(prs, "Framework 3",
    "Political and Market Barriers",
    ["⚠ Political and Incumbent Resistance — Open challenge",
     "",
     "Livestock stakeholders may resist framing their sector as a climate problem.",
     "",
     "Strategic Reframing:",
     "",
     "❌ NOT: \"Your industry is a climate problem that needs fixing through regulation.\"",
     "",
     "✓ BUT: \"Here's an investment opportunity that strengthens competitiveness",
     "         while delivering climate outcomes.\""],
    RgbColor(0x7F, 0x77, 0xDD)
)

# SLIDE 14: Co-Benefits
add_content_slide(prs, "Framework 4",
    "Co-Benefits Analysis",
    ["In the current political environment, co-benefits are often the critical path to adoption.",
     "",
     "💰 Rural Economic Development",
     "   AMC payments create new revenue streams for producers",
     "",
     "🛡️ Food System Resilience",
     "   Reduces regulatory risk, strengthens social license",
     "",
     "🌍 Trade Competitiveness",
     "   Premium positioning in markets with carbon border adjustments",
     "",
     "🔬 Innovation Spillovers",
     "   IP and capacity for other livestock challenges",
     "",
     "🤝 Global Health Equity",
     "   Scaled to LMICs, reduces climate damages for vulnerable populations"],
    RgbColor(0x7F, 0x77, 0xDD)
)

# SLIDE 15: Data Section
add_section_slide(prs, "Data & Analysis",
    "The Case in Charts",
    "Interactive data exploration supports the narrative with rigorous economics"
)

# SLIDE 16: Risk Transformation
add_content_slide(prs, "Data Analysis",
    "Innovation stalls when the downside is unbounded",
    ["Without guaranteed demand, R&D firms face outcomes that include deep losses.",
     "",
     "Firm Payoff Distribution:",
     "",
     "Without AMC: Wide distribution with significant downside risk",
     "                 ↓",
     "With AMC: Distribution shifted right, left tail truncated",
     "                 ↓",
     "Result: Losses below zero are eliminated. Investment becomes viable."],
    RgbColor(0x37, 0x8A, 0xDD)
)

# SLIDE 17: Adoption & Impact
add_stat_slide(prs, "Data Analysis",
    "Cumulative impact diverges rapidly",
    "AMC moves the adoption curve to the right and accelerates cumulative emissions reduction.",
    [("150 Mt", "Cumulative abatement by 2035 without AMC", GRAY),
     ("292 Mt", "Cumulative abatement by 2035 with AMC", TEAL),
     ("+95%", "More abatement with AMC", AMBER)],
    "",
    RgbColor(0x37, 0x8A, 0xDD)
)

# SLIDE 18: Structure Options
add_content_slide(prs, "Implementation",
    "AMC Structure Options",
    ["Governance choices set risk allocation and determine capital deployment speed.",
     "",
     "Traditional AMC — Demand-first certainty",
     "• High buyer signal before commercialization",
     "• Lower governance complexity",
     "• Longer capital deployment timeline",
     "• Proven model (pneumococcal vaccine)",
     "",
     "Hybrid AMC + Investment — Earlier capital activation",
     "• Moderate demand certainty",
     "• Faster capital timeline",
     "• Higher governance and coordination load",
     "• Can blend grant, equity, and offtake"],
    TEAL
)

# SLIDE 19: Stakeholders & Confidence
add_content_slide(prs, "Implementation",
    "Stakeholder Alignment & Confidence Tiers",
    ["Stakeholder Roles:",
     "• Innovators: R&D investment, product development, IP creation",
     "• Corporates: Demand commitment, supply chain integration",
     "• Philanthropy: Catalytic capital, risk absorption, coalition building",
     "• Governments: R&D support, regulatory pathway, verification standards",
     "",
     "Confidence Tiers:",
     "",
     "🟢 Established: Methane is addressable; AMC design principles well-understood",
     "",
     "🟡 Evolving: Vaccine efficacy timelines, dosing economics, hybrid governance",
     "",
     "🔴 Open for Engagement: Corporate commitments, government support, philanthropic capital"],
    TEAL
)

# SLIDE 20: Call to Action
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Dark green background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = TEAL_DARK
background.line.fill.background()

# Tag
tag_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(3), Inches(0.4))
tf = tag_box.text_frame
p = tf.paragraphs[0]
p.text = "CALL TO ACTION"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)

# Title
title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(11.8), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Join the Coalition"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.8), Inches(0.8))
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Led by Market Shaping Accelerator, Spark Climate, Global Methane Hub, and The Nature Conservancy"
p.font.size = Pt(16)
p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)

# Big number
num_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.8), Inches(1.5))
tf = num_box.text_frame
p = tf.paragraphs[0]
p.text = "$750M"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

label_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.5))
tf = label_box.text_frame
p = tf.paragraphs[0]
p.text = "Coalition funding target"
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)
p.alignment = PP_ALIGN.CENTER

# Stakeholders
stake_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.3), Inches(0.5))
tf = stake_box.text_frame
p = tf.paragraphs[0]
p.text = "Governments  •  Corporates  •  Investors  •  Philanthropy"
p.font.size = Pt(18)
p.font.color.rgb = RgbColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Contact
contact_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.5))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "info@lowmethanelivestock.org"
p.font.size = Pt(16)
p.font.color.rgb = RgbColor(0xA5, 0xD6, 0xA7)
p.alignment = PP_ALIGN.CENTER

# Launch window
launch_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.5))
tf = launch_box.text_frame
p = tf.paragraphs[0]
p.text = "Launch window: COP30"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Save
import io

# Write to user folder (OneDrive has write issues from Python)
final_path = r'c:\Users\delan\Methane_AMC_Slide_Deck.pptx'

# Save to BytesIO first
buffer = io.BytesIO()
prs.save(buffer)
buffer.seek(0)

# Write to file
with open(final_path, 'wb') as f:
    f.write(buffer.getvalue())

print(f"PowerPoint saved as {final_path}")
print("Note: File saved to user folder. You can move it to your preferred location.")
